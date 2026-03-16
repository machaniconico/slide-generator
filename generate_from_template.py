"""
テンプレートからスライド資料を生成するスクリプト
テンプレ内のテキストを差し替えてコンテンツを流し込む
"""

import copy
import os
import random
import subprocess
import sys
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(_SCRIPT_DIR, "templates")
OUTPUT_DIR = _SCRIPT_DIR

# フィールド別デフォルトフォントサイズ
# テンプレートのテキストボックス実測値から算出
# ボックスサイズ(内部pt) → CJK文字幅≒フォントサイズ で逆算
FIELD_FONT_SIZES = {
    "title_top": Pt(100),           # box 430×243  → 4字: 400<430 ✓
    "title_mid": Pt(36),            # box 388×51   → 8字: 288<388, h:47<51 ✓
    "title_bottom": Pt(36),         # box 366×52   → 8字: 288<366, h:47<52 ✓
    "subtitle": Pt(16),             # box 526×24   → 25字: 400<526, h:21<24 ✓
    "section_name": Pt(80),         # box 671×134  → 8字: 640<671, h:104<134 ✓
    "overview_title_top": Pt(96),   # box 264×135  → 2字: 192<264, h:125<135 ✓
    "overview_title_bottom": Pt(120),# box 917×220 → 7字: 840<917, h:156<220 ✓
    "overview_subtitle": Pt(24),    # box 347×37   → 12字: 288<347, h:31<37 ✓
    "detail": Pt(14),               # box 347×58   → 24字/行×2行=48字, h:36<58 ✓
    "body": Pt(16),                 # box 600×34   → 30字: 480<600, h:21<34 ✓
    "card_slide_title": Pt(48),     # box 677×276  → 12字: 576<677, h:62<276 ✓
    "card_title": Pt(24),           # box 290×37   → 10字: 240<290, h:31<37 ✓
    "card_desc": Pt(14),            # box 230×51   → 16字/行×2行=32字, h:36<51 ✓
    "stat_value": Pt(100),          # box 491×192  → 5hw: 275<491, h:130<192 ✓
    "stat_label": Pt(14),           # box 328×22   → 12字: 168<328, h:18<22 ✓
    "timeline_title": Pt(48),       # box 773×135  → 14字: 672<773, h:62<135 ✓
    "timeline_year": Pt(24),        # box 87×38    → 4hw: 53<87, h:31<38 ✓
    "timeline_heading": Pt(24),     # box 300×37   → 10字: 240<300, h:31<37 ✓
    "timeline_desc": Pt(14),        # box 251×51   → 15字: 210<251, h:36<51 ✓
    "closing_quote": Pt(48),        # box 994×318  → 20字: 960<994, h:62<318 ✓
}

# === フォントスタイル（10種類） ===
# Windows標準搭載フォントのみ使用
FONT_STYLES = [
    {"name": "シンプル",         "heading": "游ゴシック Medium", "body": "游ゴシック",        "heading_bold": False, "body_bold": False},
    {"name": "アーティスティック", "heading": "游明朝 Demibold",   "body": "游明朝",            "heading_bold": False, "body_bold": False},
    {"name": "コーポレート",     "heading": "BIZ UDPGothic",     "body": "BIZ UDPGothic",     "heading_bold": True,  "body_bold": False},
    {"name": "カジュアル",       "heading": "メイリオ",           "body": "メイリオ",           "heading_bold": False, "body_bold": False},
    {"name": "フォーマル",       "heading": "游明朝",            "body": "游ゴシック",        "heading_bold": False, "body_bold": False},
    {"name": "テック",           "heading": "BIZ UDPGothic",     "body": "游ゴシック",        "heading_bold": True,  "body_bold": False},
    {"name": "ナチュラル",       "heading": "游ゴシック Light",   "body": "游ゴシック Light",   "heading_bold": False, "body_bold": False},
    {"name": "ボールド",         "heading": "メイリオ",           "body": "游ゴシック Medium", "heading_bold": True,  "body_bold": False},
    {"name": "ソフト",           "heading": "UDデジタル教科書体 NK-B", "body": "UDデジタル教科書体 NK-R", "heading_bold": False, "body_bold": False},
    {"name": "ミニマル",         "heading": "游ゴシック",        "body": "游ゴシック",        "heading_bold": False, "body_bold": False},
]

# === 配色パターン（10種類） ===
# ライトテーマ前提: heading=濃い色, accent=鮮やかな差し色, body=ダークグレー
COLOR_PALETTES = [
    {"name": "ネイビー",     "heading": RGBColor(0x1A, 0x36, 0x5D), "accent": RGBColor(0x2E, 0x86, 0xDE), "body": RGBColor(0x2D, 0x2D, 0x2D)},
    {"name": "フォレスト",   "heading": RGBColor(0x1E, 0x4D, 0x3A), "accent": RGBColor(0x27, 0xAE, 0x60), "body": RGBColor(0x2D, 0x2D, 0x2D)},
    {"name": "ボルドー",     "heading": RGBColor(0x6B, 0x1D, 0x2A), "accent": RGBColor(0xC0, 0x39, 0x2B), "body": RGBColor(0x33, 0x33, 0x33)},
    {"name": "スレート",     "heading": RGBColor(0x2C, 0x3E, 0x50), "accent": RGBColor(0x5D, 0x6D, 0x7E), "body": RGBColor(0x44, 0x44, 0x44)},
    {"name": "インディゴ",   "heading": RGBColor(0x2E, 0x1A, 0x6B), "accent": RGBColor(0x8E, 0x44, 0xAD), "body": RGBColor(0x2D, 0x2D, 0x2D)},
    {"name": "サンセット",   "heading": RGBColor(0x7D, 0x3C, 0x08), "accent": RGBColor(0xE6, 0x7E, 0x22), "body": RGBColor(0x33, 0x33, 0x33)},
    {"name": "オーシャン",   "heading": RGBColor(0x0A, 0x3D, 0x62), "accent": RGBColor(0x00, 0x97, 0xE6), "body": RGBColor(0x2D, 0x2D, 0x2D)},
    {"name": "モノクロ",     "heading": RGBColor(0x1A, 0x1A, 0x1A), "accent": RGBColor(0x55, 0x55, 0x55), "body": RGBColor(0x33, 0x33, 0x33)},
    {"name": "テラコッタ",   "heading": RGBColor(0x6E, 0x3B, 0x2A), "accent": RGBColor(0xD4, 0xA5, 0x74), "body": RGBColor(0x3D, 0x30, 0x28)},
    {"name": "ミント",       "heading": RGBColor(0x0E, 0x4D, 0x45), "accent": RGBColor(0x1A, 0xBC, 0x9C), "body": RGBColor(0x2D, 0x2D, 0x2D)},
]

# フィールドタイプ → ロール（フォント・カラー両方で参照）
FIELD_STYLE_ROLE = {
    "title_top": "heading",
    "title_mid": "heading",
    "title_bottom": "heading",
    "subtitle": "body",
    "section_name": "heading",
    "overview_title_top": "accent",
    "overview_title_bottom": "heading",
    "overview_subtitle": "accent",
    "detail": "body",
    "body": "body",
    "card_slide_title": "heading",
    "card_title": "accent",
    "card_desc": "body",
    "stat_value": "accent",
    "stat_label": "body",
    "timeline_title": "heading",
    "timeline_year": "accent",
    "timeline_heading": "accent",
    "timeline_desc": "body",
    "closing_quote": "heading",
}

# アクティブなスタイル（generate_presentation で設定）
_active_font_style = None
_active_color_palette = None


def _is_light_color(rgb_color):
    """色が明るい（ダーク背景上のテキスト）かどうか判定"""
    hex_str = str(rgb_color)
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return (r + g + b) > 400


def _set_ea_font(run, font_name):
    """East Asianフォント（日本語フォント）をXMLレベルで設定"""
    from pptx.oxml.ns import qn
    from lxml import etree
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)


def _get_style_font(field_type, original_font, original_bold):
    """フィールドタイプとアクティブスタイルからフォント・太字を決定"""
    if _active_font_style is None or field_type is None:
        return original_font, original_bold

    role = FIELD_STYLE_ROLE.get(field_type)
    if role in ("heading", "accent"):
        return _active_font_style["heading"], _active_font_style.get("heading_bold", original_bold)
    elif role == "body":
        return _active_font_style["body"], _active_font_style.get("body_bold", original_bold)

    return original_font, original_bold


def _get_style_color(field_type, original_color):
    """フィールドタイプとアクティブ配色から使用する色を決定"""
    if _active_color_palette is None or field_type is None:
        return original_color

    # 元の色が明るい → ダーク背景上のテキストなので上書きしない
    if original_color and _is_light_color(original_color):
        return original_color

    role = FIELD_STYLE_ROLE.get(field_type)
    if role and role in _active_color_palette:
        return _active_color_palette[role]

    return original_color


def list_templates():
    """利用可能なテンプレート一覧"""
    templates = []
    for f in os.listdir(TEMPLATES_DIR):
        if f.endswith(".pptx"):
            templates.append(f)
    return templates


def pick_template(name=None):
    """テンプレートを選択（指定 or ランダム）"""
    templates = list_templates()
    if not templates:
        raise FileNotFoundError("templates/ にテンプレートがありません")

    if name:
        matches = [t for t in templates if name.lower() in t.lower()]
        if matches:
            return os.path.join(TEMPLATES_DIR, matches[0])
        raise FileNotFoundError(f"'{name}' に一致するテンプレートが見つかりません")

    return os.path.join(TEMPLATES_DIR, random.choice(templates))


def find_textboxes(slide):
    """スライド内のTextBoxをname→shapeのdictで返す"""
    result = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            result[shape.name] = shape
    return result


def is_cjk(ch):
    """文字がCJK（全角）か判定"""
    cp = ord(ch)
    return ((0x3000 <= cp <= 0x9FFF) or (0xF900 <= cp <= 0xFAFF)
            or (0xFF00 <= cp <= 0xFFEF) or (0x2E80 <= cp <= 0x2FFF))


def estimate_text_width(text, font_pt):
    """テキストの推定幅（pt）を返す"""
    width = 0
    for ch in text:
        if is_cjk(ch):
            width += font_pt * 1.0   # 全角: ほぼフォントサイズと同じ幅
        elif ch in (' ', '\t'):
            width += font_pt * 0.3
        elif ch in ('.,;:!?'):
            width += font_pt * 0.35
        elif ch.isupper():
            width += font_pt * 0.7   # 大文字は幅広め
        else:
            width += font_pt * 0.55  # 小文字・数字
    return width


def get_inner_size(shape):
    """テキストボックスのマージンを考慮した内部サイズ（pt）を返す"""
    tf = shape.text_frame
    # デフォルトマージン: 左右約7.2pt, 上下約3.6pt
    margin_l = (tf.margin_left or 91440) / 12700   # EMU→pt
    margin_r = (tf.margin_right or 91440) / 12700
    margin_t = (tf.margin_top or 45720) / 12700
    margin_b = (tf.margin_bottom or 45720) / 12700

    inner_w = shape.width / 12700 - margin_l - margin_r
    inner_h = shape.height / 12700 - margin_t - margin_b
    return max(inner_w, 10), max(inner_h, 10)


def calc_fit_size(shape, text, original_size):
    """テキストがシェイプに収まるフォントサイズを計算"""
    if not original_size or not shape.width or not text:
        return original_size

    inner_w, inner_h = get_inner_size(shape)
    font_pt = original_size / 12700

    # 現在のサイズでのテキスト幅
    text_width = estimate_text_width(text, font_pt)

    if text_width <= inner_w:
        # 1行に収まる → そのまま
        return original_size

    # 折り返しが必要 → 行数チェック
    line_height = font_pt * 1.3
    max_lines = max(1, int(inner_h / line_height))
    lines_needed = text_width / inner_w

    if lines_needed <= max_lines:
        # 折り返せば収まる → そのまま
        return original_size

    # 収まらない → 二分探索でフィットするサイズを見つける
    lo, hi = 6, font_pt
    best = lo
    for _ in range(20):
        mid = (lo + hi) / 2
        w = estimate_text_width(text, mid)
        lh = mid * 1.3
        ml = max(1, int(inner_h / lh))
        ln = w / inner_w

        if ln <= ml:
            best = mid
            lo = mid + 0.5
        else:
            hi = mid - 0.5

    return Pt(max(6, int(best)))


def set_autofit_text(shape):
    """テキストフレームにautofit（テキスト縮小）属性をXMLレベルで設定"""
    from pptx.oxml.ns import qn
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        # 既存のfit設定を削除
        for child in list(bodyPr):
            if child.tag in (qn('a:normAutofit'), qn('a:spAutoFit'), qn('a:noAutofit')):
                bodyPr.remove(child)
        # normAutofit を追加（PowerPointの「テキストに合わせて縮小」）
        from lxml import etree
        autofit = etree.SubElement(bodyPr, qn('a:normAutofit'))
        autofit.set('fontScale', '100000')  # 初期値100%、PowerPointが自動調整


def replace_text_keep_format(shape, new_text, field_type=None):
    """テキストを差し替え（フォーマット維持 + 自動サイズ調整 + autofit）"""
    tf = shape.text_frame
    tf.word_wrap = True
    if not tf.paragraphs:
        return

    # 最初のパラグラフの最初のランのフォーマットを取得
    first_para = tf.paragraphs[0]
    if first_para.runs:
        ref_run = first_para.runs[0]
        font_name = ref_run.font.name
        font_size = ref_run.font.size
        font_bold = ref_run.font.bold
        font_color = ref_run.font.color.rgb if ref_run.font.color and ref_run.font.color.rgb else None
    else:
        font_name = None
        font_size = None
        font_bold = None
        font_color = None

    # フィールド別デフォルトサイズで上書き
    if field_type and field_type in FIELD_FONT_SIZES:
        font_size = FIELD_FONT_SIZES[field_type]

    # フォントサイズを計算
    adjusted_size = calc_fit_size(shape, new_text, font_size)

    # 全パラグラフをクリア
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    # スタイル上書き（フォント・太字・カラー）
    font_name, font_bold = _get_style_font(field_type, font_name, font_bold)
    final_color = _get_style_color(field_type, font_color)

    # 最初のパラグラフにテキストをセット
    first_para.clear()
    run = first_para.add_run()
    run.text = new_text
    if font_name:
        run.font.name = font_name
        _set_ea_font(run, font_name)
    if adjusted_size:
        run.font.size = adjusted_size
    if font_bold is not None:
        run.font.bold = font_bold
    if final_color:
        run.font.color.rgb = final_color

    # PowerPointの自動縮小もフォールバックとして有効化
    set_autofit_text(shape)


def replace_text_multiline(shape, lines):
    """複数行テキストを差し替え（各行をパラグラフとして追加）"""
    tf = shape.text_frame

    # 最初のパラグラフからフォーマットを取得
    first_para = tf.paragraphs[0]
    if first_para.runs:
        ref_run = first_para.runs[0]
        font_name = ref_run.font.name
        font_size = ref_run.font.size
        font_bold = ref_run.font.bold
        font_color = ref_run.font.color.rgb if ref_run.font.color and ref_run.font.color.rgb else None
        alignment = first_para.alignment
    else:
        font_name = None
        font_size = None
        font_bold = None
        font_color = None
        alignment = None

    # 全パラグラフをクリア
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)

    # 最初の行
    first_para.clear()
    run = first_para.add_run()
    run.text = lines[0]
    if font_name:
        run.font.name = font_name
    if font_size:
        run.font.size = font_size
    if font_bold is not None:
        run.font.bold = font_bold
    if font_color:
        run.font.color.rgb = font_color

    # 残りの行を追加
    from pptx.oxml.ns import qn
    from copy import deepcopy
    for line in lines[1:]:
        new_p = deepcopy(first_para._p)
        # テキストを設定
        for r in new_p.findall(qn('a:r')):
            for t in r.findall(qn('a:t')):
                t.text = line
        tf._txBody.append(new_p)


def delete_slides(prs, slide_indices_to_keep):
    """指定したインデックスのスライドだけ残して他を削除"""
    # 逆順で削除
    all_indices = set(range(len(prs.slides)))
    to_remove = sorted(all_indices - set(slide_indices_to_keep), reverse=True)

    for idx in to_remove:
        rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldId = prs.slides._sldIdLst[idx]
        prs.slides._sldIdLst.remove(sldId)


def generate_presentation(template_path, content, output_name, style_name=None, color_name=None):
    """
    content = {
        "title": "メインタイトル",
        "subtitle": "サブタイトル",
        "sections": [
            {
                "number": "01",
                "name": "セクション名",
                "slides": [
                    {"type": "content", "title": "...", "body": "..."},
                    {"type": "cards", "items": [("title","body"), ...]},
                    {"type": "stats", "items": [("number","label"), ...]},
                ]
            },
        ],
        "closing_quote": "クロージングメッセージ"
    }
    style_name: フォントスタイル名（省略でランダム）
    color_name: 配色パターン名（省略でランダム）
    """
    global _active_font_style, _active_color_palette

    if style_name:
        matches = [s for s in FONT_STYLES if s["name"] == style_name]
        _active_font_style = matches[0] if matches else random.choice(FONT_STYLES)
    else:
        _active_font_style = random.choice(FONT_STYLES)

    if color_name:
        matches = [c for c in COLOR_PALETTES if c["name"] == color_name]
        _active_color_palette = matches[0] if matches else random.choice(COLOR_PALETTES)
    else:
        _active_color_palette = random.choice(COLOR_PALETTES)

    print(f"スタイル: {_active_font_style['name']} × {_active_color_palette['name']}")

    prs = Presentation(template_path)

    # テンプレートのスライドを解析してテキスト差し替え
    slides = list(prs.slides)

    # --- スライド1: 表紙 ---
    s1_texts = find_textboxes(slides[0])
    for name, shape in s1_texts.items():
        text = shape.text_frame.text.strip()
        if text == "SEO":
            replace_text_keep_format(shape, content.get("title_top", ""), "title_top")
        elif text == "ISOMETRIC":
            replace_text_keep_format(shape, content.get("title_mid", ""), "title_mid")
        elif text == "STRATEGY":
            replace_text_keep_format(shape, content.get("title_bottom", ""), "title_bottom")

    # --- スライド2: 目次 ---
    s2_texts = find_textboxes(slides[1])
    sections = content.get("sections", [])
    section_names = {
        "INTRODUCTION": sections[0]["name"] if len(sections) > 0 else "",
        "OUR PROJECTS": sections[1]["name"] if len(sections) > 1 else "",
        "ABOUT US": sections[2]["name"] if len(sections) > 2 else "",
    }
    for name, shape in s2_texts.items():
        text = shape.text_frame.text.strip()
        if text in section_names:
            replace_text_keep_format(shape, section_names[text], "section_name")

    # --- スライド3: セクション1区切り ---
    s3_texts = find_textboxes(slides[2])
    for name, shape in s3_texts.items():
        text = shape.text_frame.text.strip()
        if text == "INTRODUCTION":
            replace_text_keep_format(shape, sections[0]["name"] if sections else "", "section_name")
        elif "SEO STRATEGY" in text:
            replace_text_keep_format(shape, content.get("subtitle", ""), "subtitle")
        elif "Elaborate" in text:
            body = sections[0]["slides"][0]["body"] if sections and sections[0]["slides"] else ""
            replace_text_keep_format(shape, body, "body")

    # --- スライド5: 概要コンテンツ ---
    s5_texts = find_textboxes(slides[4])
    for name, shape in s5_texts.items():
        text = shape.text_frame.text.strip()
        if text == "OUR":
            replace_text_keep_format(shape, sections[0]["slides"][0].get("title_top", ""), "overview_title_top")
        elif text == "COMPANY":
            replace_text_keep_format(shape, sections[0]["slides"][0].get("title_bottom", ""), "overview_title_bottom")
        elif "What we do" in text:
            replace_text_keep_format(shape, sections[0]["slides"][0].get("subtitle", ""), "overview_subtitle")
        elif "Briefly elaborate" in text:
            replace_text_keep_format(shape, sections[0]["slides"][0].get("detail", ""), "detail")

    # --- スライド7: セクション2区切り ---
    s7_texts = find_textboxes(slides[6])
    for name, shape in s7_texts.items():
        text = shape.text_frame.text.strip()
        if text == "OUR PROJECTS":
            replace_text_keep_format(shape, sections[1]["name"] if len(sections) > 1 else "", "section_name")
        elif "SEO STRATEGY" in text:
            replace_text_keep_format(shape, content.get("subtitle", ""), "subtitle")
        elif "Elaborate" in text:
            body = sections[1]["slides"][0]["body"] if len(sections) > 1 and sections[1]["slides"] else ""
            replace_text_keep_format(shape, body, "body")

    # --- スライド8: サービス/カード ---
    s8_texts = find_textboxes(slides[7])
    cards = sections[1]["slides"][0].get("cards", []) if len(sections) > 1 and sections[1]["slides"] else []
    card_titles = ["Service One", "Service Two", "Service Three", "Service Four"]
    elaborate_count = 0
    for name, shape in s8_texts.items():
        text = shape.text_frame.text.strip()
        if text == "WHAT WE OFFER":
            replace_text_keep_format(shape, sections[1]["slides"][0].get("title", ""), "card_slide_title")
        elif text in card_titles:
            idx = card_titles.index(text)
            if idx < len(cards):
                replace_text_keep_format(shape, cards[idx][0], "card_title")
        elif "Elaborate" in text:
            if elaborate_count < len(cards):
                replace_text_keep_format(shape, cards[elaborate_count][1], "card_desc")
            elaborate_count += 1

    # --- スライド10: 統計 ---
    # テンプレの配置順: "5.000"(y=59), "20K +"(y=300), "X 1.5"(y=541)
    # ラベル配置順:     y=228,          y=469,           y=710
    stats = sections[1]["slides"][1].get("stats", []) if len(sections) > 1 and len(sections[1]["slides"]) > 1 else []
    stat_value_order = ["5.000", "20K +", "X 1.5"]  # 上から順

    # 位置でソートしたラベルシェイプを集める
    stat_value_shapes = []
    label_shapes = []
    for shape in slides[9].shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text in stat_value_order:
            stat_value_shapes.append((shape.top, shape, stat_value_order.index(text)))
        elif "Elaborate" in text or "featured statistic" in text:
            label_shapes.append((shape.top, shape))

    stat_value_shapes.sort(key=lambda x: x[0])  # 上から順
    label_shapes.sort(key=lambda x: x[0])

    for i, (_, shape, _) in enumerate(stat_value_shapes):
        if i < len(stats):
            replace_text_keep_format(shape, stats[i][0], "stat_value")
    for i, (_, shape) in enumerate(label_shapes):
        if i < len(stats):
            replace_text_keep_format(shape, stats[i][1], "stat_label")

    # --- スライド13: セクション3区切り ---
    s13_texts = find_textboxes(slides[12])
    for name, shape in s13_texts.items():
        text = shape.text_frame.text.strip()
        if text == "ABOUT US":
            replace_text_keep_format(shape, sections[2]["name"] if len(sections) > 2 else "", "section_name")
        elif "SEO STRATEGY" in text:
            replace_text_keep_format(shape, content.get("subtitle", ""), "subtitle")
        elif "Elaborate" in text:
            body = sections[2]["slides"][0]["body"] if len(sections) > 2 and sections[2]["slides"] else ""
            replace_text_keep_format(shape, body, "body")

    # --- スライド15: タイムライン ---
    s15_texts = find_textboxes(slides[14])
    timeline = sections[2]["slides"][0].get("timeline", []) if len(sections) > 2 and sections[2]["slides"] else []
    year_labels = ["2005", "2015", "2025"]
    point_count = 0
    elaborate_count = 0
    for name, shape in s15_texts.items():
        text = shape.text_frame.text.strip()
        if text == "OUR HISTORY":
            replace_text_keep_format(shape, sections[2]["slides"][0].get("title", ""), "timeline_title")
        elif text in year_labels:
            idx = year_labels.index(text)
            if idx < len(timeline):
                replace_text_keep_format(shape, timeline[idx][0], "timeline_year")
        elif "Add a main point" in text:
            if point_count < len(timeline):
                replace_text_keep_format(shape, timeline[point_count][1], "timeline_heading")
            point_count += 1
        elif "Elaborate" in text:
            if elaborate_count < len(timeline):
                replace_text_keep_format(shape, timeline[elaborate_count][2], "timeline_desc")
            elaborate_count += 1

    # --- スライド17: クロージング ---
    s17_texts = find_textboxes(slides[16])
    for name, shape in s17_texts.items():
        text = shape.text_frame.text.strip()
        if "original statement" in text or "inspiring quote" in text:
            replace_text_keep_format(shape, content.get("closing_quote", "Thank you"), "closing_quote")

    # 使用するスライドだけ残す
    keep = [0, 1, 2, 4, 6, 7, 9, 12, 14, 16]
    delete_slides(prs, keep)

    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)
    print(f"生成完了: {output_path}")
    return output_path


def _find_libreoffice():
    """LibreOfficeの実行パスを探す"""
    candidates = [
        "libreoffice",
        "C:/Program Files/LibreOffice/program/soffice.exe",
        "C:/Program Files (x86)/LibreOffice/program/soffice.exe",
    ]
    for path in candidates:
        try:
            subprocess.run([path, "--version"], capture_output=True, timeout=10)
            return path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


def convert_to_pdf(pptx_path, pdf_path=None):
    """PPTXをPDFに変換（PowerPoint → LibreOffice の順で試行）"""
    if pdf_path is None:
        pdf_path = os.path.splitext(pptx_path)[0] + '.pdf'

    pptx_abs = os.path.abspath(pptx_path).replace('/', '\\')
    pdf_abs = os.path.abspath(pdf_path).replace('/', '\\')

    # 方法1: PowerPointデスクトップ版（フォント完全一致）
    try:
        ps_cmd = '; '.join([
            '$ErrorActionPreference = "Stop"',
            '$ppt = New-Object -ComObject PowerPoint.Application',
            f'$pres = $ppt.Presentations.Open("{pptx_abs}")',
            f'$pres.SaveAs("{pdf_abs}", 32)',
            '$pres.Close()',
            '$ppt.Quit()',
        ])
        subprocess.run(
            ['powershell.exe', '-NoProfile', '-Command', ps_cmd],
            check=True, capture_output=True, text=True, timeout=120
        )
        print(f"PDF出力（PowerPoint）: {pdf_path}")
        return pdf_path
    except (subprocess.CalledProcessError, FileNotFoundError):
        pass

    # 方法2: LibreOffice（フォントが近い代替に置換される場合あり）
    lo_path = _find_libreoffice()
    if lo_path:
        try:
            out_dir = os.path.dirname(pptx_abs) or '.'
            subprocess.run(
                [lo_path, '--headless', '--convert-to', 'pdf', '--outdir', out_dir, pptx_abs],
                check=True, capture_output=True, text=True, timeout=120
            )
            print(f"PDF出力（LibreOffice）: {pdf_path}")
            print("  ※ フォントが一部代替される場合があります")
            return pdf_path
        except subprocess.CalledProcessError as e:
            error_msg = e.stderr.strip() if e.stderr else str(e)
            raise RuntimeError(f"LibreOffice PDF変換エラー: {error_msg}")

    # どちらもない
    raise RuntimeError(
        "PDF変換に必要なソフトが見つかりません。\n"
        "以下のいずれかをインストールしてください:\n"
        "  1. PowerPointデスクトップ版（推奨・フォント完全一致）\n"
        "  2. LibreOffice（無料・フォントが近似になる場合あり）\n"
        "     https://www.libreoffice.org/download/"
    )


# === サンプルコンテンツ（生成AIトレンド） ===
SAMPLE_CONTENT = {
    "title_top": "生成AI",
    "title_mid": "2025年",
    "title_bottom": "トレンド",
    "subtitle": "主要な技術動向と今後の展望",

    "sections": [
        {
            "number": "01",
            "name": "マルチモーダル",
            "slides": [
                {
                    "type": "overview",
                    "title_top": "AI",
                    "title_bottom": "マルチモーダル",
                    "subtitle": "統合理解の時代へ",
                    "detail": "テキスト・画像・音声・動画を統合処理するモデルが主流に。リアルタイム対話や動画生成が実用化。",
                    "body": "GPT-4o、Gemini、Claude等が対応を強化。2025年は標準機能として定着。",
                },
            ],
        },
        {
            "number": "02",
            "name": "AIエージェント",

            "slides": [
                {
                    "type": "cards",
                    "title": "AIエージェントの主要カテゴリ",
                    "body": "自律型AIエージェントが急速に普及。開発・業務・リサーチで活用拡大。",
                    "cards": [
                        ("コーディング", "PR作成からバグ修正まで自動対応"),
                        ("業務自動化", "定型業務の自律実行が加速"),
                        ("リサーチ", "情報収集と統合を自動化"),
                        ("マルチAgent", "複数AIが協調して問題解決"),
                    ],
                },
                {
                    "type": "stats",
                    "stats": [
                        ("300%", "市場の前年比成長率"),
                        ("x 2.5", "生産性向上の倍率"),
                        ("85%", "大企業の導入率"),
                    ],
                },
            ],
        },
        {
            "number": "03",
            "name": "オープンソース",
            "slides": [
                {
                    "type": "timeline",
                    "title": "オープンソースAIの進化",
                    "timeline": [
                        ("2023", "Llama 2公開", "Meta、商用利用を許可"),
                        ("2024", "性能の急接近", "GPT-4に迫る性能を達成"),
                        ("2025", "実用化の加速", "企業導入が加速"),
                    ],
                    "body": "APIの代替から独自の強みを持つ選択肢へ進化。",
                },
            ],
        },
    ],

    "closing_quote": "AIは道具を超え、パートナーになりつつある",
}


if __name__ == "__main__":
    # --pdf フラグを検出して除去
    pdf_mode = "--pdf" in sys.argv
    argv = [a for a in sys.argv[1:] if a != "--pdf"]

    # 引数: [テンプレ名] [スタイル名] [配色名]（省略でランダム）
    template_name = argv[0] if len(argv) > 0 else None
    style_name = argv[1] if len(argv) > 1 else None
    color_name = argv[2] if len(argv) > 2 else None

    if template_name == "list":
        print("利用可能なテンプレート:")
        for t in list_templates():
            print(f"  - {t}")
        print(f"\nフォントスタイル（{len(FONT_STYLES)}種）:")
        for s in FONT_STYLES:
            print(f"  - {s['name']}")
        print(f"\n配色パターン（{len(COLOR_PALETTES)}種）:")
        for c in COLOR_PALETTES:
            print(f"  - {c['name']}")
        print(f"\n組み合わせ: {len(FONT_STYLES) * len(COLOR_PALETTES)}通り")
        sys.exit(0)

    template_path = pick_template(template_name)
    print(f"テンプレート: {os.path.basename(template_path)}")

    output_pptx = "output_presentation.pptx"
    output_path = generate_presentation(template_path, SAMPLE_CONTENT, output_pptx, style_name, color_name)

    if pdf_mode:
        convert_to_pdf(output_path)
